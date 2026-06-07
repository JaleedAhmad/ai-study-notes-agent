import sys
import os
import io

# Add parent to path for imports
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))

from src.utils.extractor import extract_universal_text

def create_mock_file(name, content_bytes):
    f = io.BytesIO(content_bytes)
    f.name = name
    return f

def test_markdown():
    print("Testing Markdown extraction...")
    content = b"# Hello Markdown\nThis is a test."
    file = create_mock_file("test.md", content)
    text = extract_universal_text(file)
    print(f"Extracted: {text.strip()}")
    assert "Hello Markdown" in text

def test_docx():
    print("Testing DOCX extraction...")
    import docx
    doc = docx.Document()
    doc.add_paragraph("Hello from DOCX")
    f = io.BytesIO()
    doc.save(f)
    file = create_mock_file("test.docx", f.getvalue())
    text = extract_universal_text(file)
    print(f"Extracted: {text.strip()}")
    assert "Hello from DOCX" in text

def test_pptx():
    print("Testing PPTX extraction...")
    import pptx
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Hello from PPTX"
    f = io.BytesIO()
    prs.save(f)
    file = create_mock_file("test.pptx", f.getvalue())
    text = extract_universal_text(file)
    print(f"Extracted: {text.strip()}")
    assert "Hello from PPTX" in text

if __name__ == "__main__":
    test_markdown()
    test_docx()
    test_pptx()
    print("All tests passed!")
