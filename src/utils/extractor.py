import io
import docx
import pptx
from pypdf import PdfReader
from PIL import Image
from pdf2image import convert_from_bytes
from ..core.vision_client import extract_text_from_image

def extract_universal_text(uploaded_file):
    """
    Extracts text from uploaded files universally supporting PDF, DOCX, PPTX, TXT, and MD.
    Expects a Streamlit UploadedFile or file-like object with a `.name` attribute.
    """
    try:
        filename = uploaded_file.name.lower()
        
        if hasattr(uploaded_file, "seek"):
            uploaded_file.seek(0)
            
        if filename.endswith(".pdf"):
            reader = PdfReader(uploaded_file)
            text = ""
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
            
            # Smart Fallback for scanned PDFs
            num_pages = len(reader.pages)
            if num_pages > 0 and len(text.strip()) / num_pages < 50:
                print("LOG: [Extractor] -> PDF text suspiciously short, falling back to Vision OCR via pdf2image...")
                if hasattr(uploaded_file, "seek"):
                    uploaded_file.seek(0)
                pdf_bytes = uploaded_file.read()
                images = convert_from_bytes(pdf_bytes)
                text = ""
                for idx, img in enumerate(images):
                    print(f"LOG: [Extractor] -> OCR on PDF page {idx+1}/{len(images)}...")
                    page_text = extract_text_from_image(img)
                    text += f"\\n--- Page {idx+1} ---\\n{page_text}\\n"
                    
            return text            
        elif filename.endswith(".docx"):
            doc = docx.Document(uploaded_file)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
            
        elif filename.endswith(".pptx"):
            prs = pptx.Presentation(uploaded_file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        text += shape.text + "\n"
            return text
            
        elif filename.endswith(".txt") or filename.endswith(".md"):
            content = uploaded_file.read()
            if isinstance(content, bytes):
                return content.decode("utf-8", errors="replace")
            return str(content)
            
        elif filename.endswith((".png", ".jpg", ".jpeg")):
            img = Image.open(uploaded_file)
            if img.mode != "RGB":
                img = img.convert("RGB")
            text = extract_text_from_image(img)
            return text
            
        else:
            return ""
            
    except Exception as e:
        print(f"Extraction error: {e}")
        return ""
